from typing import Annotated, Tuple, Optional
from urllib.parse import urlparse, urlunparse
import logging
import sys
import time
import io
import os
import re
import hashlib
import tempfile
import mimetypes
from pathlib import Path

import markdownify
import docx
import PyPDF2
from pptx import Presentation
import readabilipy.simple_json
from mcp.shared.exceptions import McpError
from mcp.server import Server
from mcp.server.stdio import stdio_server
from mcp.types import (
    ErrorData,
    GetPromptResult,
    Prompt,
    PromptArgument,
    PromptMessage,
    TextContent,
    Tool,
    INVALID_PARAMS,
    INTERNAL_ERROR,
)
# Removed robots.txt related import
from pydantic import BaseModel, Field, AnyUrl

# Additional imports for browser automation and OCR
import requests
from bs4 import BeautifulSoup
import undetected_chromedriver as uc
from selenium.webdriver.common.by import By
import numpy as np
from PIL import Image
import pytesseract

# Set up logger
logger = logging.getLogger("mcp-fetch")

DEFAULT_USER_AGENT_AUTONOMOUS = "ModelContextProtocol/1.0 (Autonomous; +https://github.com/modelcontextprotocol/servers)"
DEFAULT_USER_AGENT_MANUAL = "ModelContextProtocol/1.0 (User-Specified; +https://github.com/modelcontextprotocol/servers)"

# Try to import layoutparser, but don't fail if it's not available
try:
    import layoutparser as lp
    LAYOUT_PARSER_AVAILABLE = True
except ModuleNotFoundError:
    lp = None
    LAYOUT_PARSER_AVAILABLE = False
    logger.warning("layoutparser not installed. Layout detection is disabled.")


def _cleanup_extracted_text(text: str) -> str:
    """Clean up extracted text by stripping blank lines and extra whitespace."""
    lines = [line.strip() for line in text.splitlines()]
    # Remove empty lines and consecutive duplicates
    clean_lines = []
    prev_line = None
    for line in lines:
        if line and line != prev_line:
            clean_lines.append(line)
            prev_line = line
    return "\n".join(clean_lines)


def extract_content_from_html(html: str) -> str:
    """Extract and convert HTML content to Markdown format.

    Args:
        html: Raw HTML content to process

    Returns:
        Simplified markdown version of the content
    """
    logger.debug("Converting HTML content to Markdown (content length: %d bytes)", len(html))
    start_time = time.time()
    ret = readabilipy.simple_json.simple_json_from_html_string(
        html, use_readability=True
    )
    if not ret["content"]:
        logger.warning("HTML simplification failed - no content extracted")
        return "<error>Page failed to be simplified from HTML</error>"
    content = markdownify.markdownify(
        ret["content"],
        heading_style=markdownify.ATX,
    )
    logger.debug(
        "HTML conversion completed in %.2f seconds (output length: %d bytes)",
        time.time() - start_time, 
        len(content)
    )
    return content


def get_robots_txt_url(url: str) -> str:
    """Get the robots.txt URL for a given website URL.

    Args:
        url: Website URL to get robots.txt for

    Returns:
        URL of the robots.txt file
    """
    # Parse the URL into components
    parsed = urlparse(url)

    # Reconstruct the base URL with just scheme, netloc, and /robots.txt path
    robots_url = urlunparse((parsed.scheme, parsed.netloc, "/robots.txt", "", "", ""))
    
    logger.debug("Generated robots.txt URL: %s from base URL: %s", robots_url, url)
    return robots_url

# Removed check_may_autonomously_fetch_url function since robots.txt check is no longer needed.


def _capture_screenshot(url: str) -> Tuple[bytes, str]:
    """
    Capture a full-page screenshot and page source using undetected-chromedriver.
    It also clicks common cookie banners before capturing.
    """
    driver = None
    try:
        logger.debug("Launching undetected-chromedriver for screenshot capture.")
        options = uc.ChromeOptions()
        options.add_argument("--headless")
        options.add_argument("--no-sandbox")
        options.add_argument("--disable-dev-shm-usage")
        options.add_argument("--disable-gpu")
        options.add_argument("--window-size=1920,1080")
        
        driver = uc.Chrome(options=options)
        driver.set_page_load_timeout(30)
        
        logger.debug(f"Navigating to {url}")
        driver.get(url)
        time.sleep(5)  # Wait for dynamic content to load
        
        # Common cookie banner selectors
        cookie_selectors = [
            "//*[contains(translate(text(),'ABCDEFGHIJKLMNOPQRSTUVWXYZ','abcdefghijklmnopqrstuvwxyz'), 'accept')]",
            "//*[contains(translate(text(),'ABCDEFGHIJKLMNOPQRSTUVWXYZ','abcdefghijklmnopqrstuvwxyz'), 'agree')]",
            "//*[contains(translate(text(),'ABCDEFGHIJKLMNOPQRSTUVWXYZ','abcdefghijklmnopqrstuvwxyz'), 'got it')]",
            "//*[contains(translate(text(),'ABCDEFGHIJKLMNOPQRSTUVWXYZ','abcdefghijklmnopqrstuvwxyz'), 'allow')]",
            "//*[contains(translate(text(),'ABCDEFGHIJKLMNOPQRSTUVWXYZ','abcdefghijklmnopqrstuvwxyz'), 'consent')]",
            "//*[contains(translate(text(),'ABCDEFGHIJKLMNOPQRSTUVWXYZ','abcdefghijklmnopqrstuvwxyz'), 'cookies')]",
            "button#onetrust-accept-btn-handler",
            "button.accept-cookies",
            "button.cookie-accept",
            "button.cookie-consent-accept",
            "button.cookie-banner__accept",
            "#cookie-accept-btn",
            "#cookie-banner-accept",
            "#gdpr-cookie-accept",
            ".cookie-banner .accept-button",
            ".cookie-consent .accept-button"
        ]
        
        for selector in cookie_selectors:
            try:
                if selector.startswith("/"):
                    elements = driver.find_elements(By.XPATH, selector)
                    for elem in elements:
                        if elem.is_displayed():
                            try:
                                driver.execute_script("arguments[0].click();", elem)
                                logger.debug(f"Clicked cookie banner using XPath: {selector}")
                                time.sleep(1)
                                break
                            except Exception as click_e:
                                logger.debug(f"Click error for XPath {selector}: {click_e}")
                else:
                    elements = driver.find_elements(By.CSS_SELECTOR, selector)
                    for elem in elements:
                        if elem.is_displayed():
                            try:
                                driver.execute_script("arguments[0].click();", elem)
                                logger.debug(f"Clicked cookie banner using CSS: {selector}")
                                time.sleep(1)
                                break
                            except Exception as click_e:
                                logger.debug(f"Click error for CSS {selector}: {click_e}")
            except Exception as e:
                logger.debug(f"Selector {selector} not found: {e}")
        
        total_height = driver.execute_script("return Math.max(document.body.scrollHeight, document.documentElement.scrollHeight);")
        driver.set_window_size(1920, total_height)
        logger.debug(f"Adjusted window size to 1920x{total_height}")
        time.sleep(2)
        
        screenshot = driver.get_screenshot_as_png()
        logger.debug("Screenshot captured")
        
        page_source = driver.page_source
        logger.debug(f"Page source captured, length: {len(page_source)}")
        
        return screenshot, page_source
    except Exception as e:
        logger.exception(f"Screenshot capture failed: {e}")
        return None, None
    finally:
        if driver:
            try:
                driver.quit()
                logger.debug("Chrome driver closed successfully")
            except Exception as e:
                logger.error(f"Error closing Chrome driver: {e}")


def _extract_text_with_pytesseract(image_bytes: bytes) -> str:
    """Use pytesseract with layout detection to OCR the screenshot."""
    try:
        logger.debug("Extracting text with pytesseract using layout detection.")
        pil_image = Image.open(io.BytesIO(image_bytes)).convert("RGB")
        if LAYOUT_PARSER_AVAILABLE and lp is not None:
            logger.debug("layoutparser is available.")
            # Initialize Layout Parser model (using a public model from PubLayNet)
            try:
                model = lp.Detectron2LayoutModel(
                    config_path='/app/weights/config.yml',
                    model_path='/app/weights/model_final.pth',
                    extra_config=["MODEL.ROI_HEADS.SCORE_THRESH_TEST", 0.5],
                    label_map={0: "Text", 1: "Title", 2: "List", 3: "Table", 4: "Figure"}
                )
                # Convert PIL image to numpy array for layoutparser
                image_np = np.array(pil_image)
                layout = model.detect(image_np)
                logger.debug("Layoutparser successfully detected layout.")
                extracted_text = ""
                text_blocks = [block for block in layout if block.type == "Text"]
                if text_blocks:
                    for block in text_blocks:
                        x_1, y_1, x_2, y_2 = block.coordinates
                        cropped_pil = pil_image.crop((int(x_1), int(y_1), int(x_2), int(y_2)))
                        txt = pytesseract.image_to_string(cropped_pil)
                        extracted_text += txt + "\n"
                else:
                    extracted_text = pytesseract.image_to_string(pil_image)
            except Exception as e:
                logger.warning(f"Layout parser model loading failed: {e}. Using basic OCR instead.")
                extracted_text = pytesseract.image_to_string(pil_image)
        else:
            logger.warning("layoutparser is not available. Falling back to basic pytesseract OCR.")
            extracted_text = pytesseract.image_to_string(pil_image)
        return _cleanup_extracted_text(extracted_text)
    except Exception as e:
        logger.exception(f"Pytesseract extraction with layout detection failed: {e}")
        return ""


def extract_html_with_requests(url: str) -> str:
    """Extract HTML content using requests and BeautifulSoup."""
    try:
        logger.debug("Extracting HTML using requests/BeautifulSoup.")
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
        }
        response = requests.get(url, headers=headers, timeout=10)
        if response.status_code != 200:
            return ""
        soup = BeautifulSoup(response.text, "html.parser")
        for element in soup(["script", "style", "header", "footer", "nav"]):
            element.decompose()
        text = soup.get_text(separator="\n")
        return _cleanup_extracted_text(text)
    except Exception as e:
        logger.exception(f"HTML extraction failed: {e}")
        return ""


def _parse_pdf(content: bytes) -> str:
    """Parse PDF content and extract text.
    
    Args:
        content: Raw PDF file content as bytes
        
    Returns:
        Extracted text content
    """
    try:
        logger.debug("Parsing PDF content")
        pdf_file = io.BytesIO(content)
        reader = PyPDF2.PdfReader(pdf_file)
        
        extracted_text = []
        for page_num in range(len(reader.pages)):
            page = reader.pages[page_num]
            extracted_text.append(page.extract_text())
            
        text = "\n\n".join(extracted_text)
        return _cleanup_extracted_text(text)
    except Exception as e:
        logger.error(f"Failed to parse PDF: {str(e)}")
        return ""

def _parse_docx(content: bytes) -> str:
    """Parse DOCX content and extract text.
    
    Args:
        content: Raw DOCX file content as bytes
        
    Returns:
        Extracted text content
    """
    try:
        logger.debug("Parsing DOCX content")
        docx_file = io.BytesIO(content)
        doc = docx.Document(docx_file)
        
        # Extract text from paragraphs
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    paragraphs.append(row_text)
        
        text = "\n\n".join(paragraphs)
        return _cleanup_extracted_text(text)
    except Exception as e:
        logger.error(f"Failed to parse DOCX: {str(e)}")
        return ""

def _parse_pptx(content: bytes) -> str:
    """Parse PPTX content and extract text.
    
    Args:
        content: Raw PPTX file content as bytes
        
    Returns:
        Extracted text content
    """
    try:
        logger.debug("Parsing PPTX content")
        pptx_file = io.BytesIO(content)
        prs = Presentation(pptx_file)
        
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_parts = [f"Slide {slide_num}:"]
            
            # Get text from shapes (including text boxes)
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())
                    
            # Add notes if they exist
            if slide.notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                slide_parts.append(f"Notes: {slide.notes_slide.notes_text_frame.text.strip()}")
                
            if len(slide_parts) > 1:  # Only add slides that have content
                text_parts.append("\n".join(slide_parts))
        
        text = "\n\n".join(text_parts)
        return _cleanup_extracted_text(text)
    except Exception as e:
        logger.error(f"Failed to parse PPTX: {str(e)}")
        return ""

def _get_content_type(url: str, headers: Optional[dict] = None) -> str:
    """Determine content type from URL and/or headers.
    
    Args:
        url: The URL being fetched
        headers: Optional response headers containing content-type
        
    Returns:
        Detected content type (e.g., 'pdf', 'docx', 'pptx', 'html', 'text')
    """
    # First check content-type header if available
    if headers and 'content-type' in headers:
        content_type = headers['content-type'].lower()
        if 'pdf' in content_type:
            return 'pdf'
        elif 'wordprocessingml.document' in content_type or 'msword' in content_type:
            return 'docx'
        elif 'presentationml.presentation' in content_type or 'powerpoint' in content_type:
            return 'pptx'
        elif 'text/html' in content_type:
            return 'html'
        elif 'text/plain' in content_type:
            return 'text'
    
    # Fall back to URL extension
    ext = Path(urlparse(url).path).suffix.lower()
    if ext in ['.pdf']:
        return 'pdf'
    elif ext in ['.docx', '.doc']:
        return 'docx'
    elif ext in ['.pptx', '.ppt']:
        return 'pptx'
    elif ext in ['.html', '.htm']:
        return 'html'
    elif ext in ['.txt']:
        return 'text'
    
    # Default to HTML if no specific type detected
    return 'html'

def choose_best_result(results: list) -> Tuple[str, str]:
    """Choose the best result based on multiple quality criteria."""
    # Filter out empty results
    valid_results = [(name, text) for name, text in results if text and text.strip()]
    if not valid_results:
        return "none", ""
    
    def score_content(name: str, text: str) -> float:
        """Score content based on multiple quality criteria."""
        score = 0.0
        text_length = len(text)
        
        # Base score from content length (max 50 points)
        score += min(text_length / 100, 50)
        
        # Bonus points for structured content (max 20 points)
        if text.count('\n') > 0:
            paragraphs = text.count('\n\n')
            score += min(paragraphs, 20)
        
        # Penalize extremely short content
        if text_length < 100:
            score *= 0.5
        
        # Penalize likely error messages or invalid content
        error_indicators = ['<error>', 'failed to', 'error occurred', 'access denied']
        if any(indicator in text.lower() for indicator in error_indicators):
            score *= 0.1
            
        return score
    
    # Score and sort results
    scored_results = [(name, text, score_content(name, text))
                     for name, text in valid_results]
    sorted_results = sorted(scored_results, key=lambda x: x[2], reverse=True)
    
    # Log scores for debugging
    logger.debug("Content scoring results:")
    for name, _, score in sorted_results:
        logger.debug(f"{name}: {score:.2f}")
    
    return sorted_results[0][0], sorted_results[0][1]


async def fetch_url_with_multiple_methods(url: str, user_agent: str) -> Tuple[str, str]:
    """
    Fetch the URL using multiple methods and return the best result.
    Supports various file formats including PDF, DOCX, PPTX, and HTML.
    """
    extracted_texts = {}
    content_type = None
    raw_content = None

    # Try to fetch content first with httpx to detect content type
    try:
        from httpx import AsyncClient, HTTPError
        async with AsyncClient() as client:
            response = await client.get(
                url,
                follow_redirects=True,
                headers={"User-Agent": user_agent},
                timeout=30,
            )
            if response.status_code < 400:
                content_type = _get_content_type(url, dict(response.headers))
                raw_content = response.content  # Get raw bytes for file parsing
                logger.info(f"Detected content type: {content_type}")
    except Exception as e:
        logger.error("Failed to fetch with httpx: %s", str(e))

    # Handle document formats (PDF, DOCX, PPTX, TXT)
    if content_type in ['pdf', 'docx', 'pptx', 'text']:
        logger.info(f"Processing as {content_type.upper()} document")
        
        # Try parsing with httpx response first
        if raw_content:
            if content_type == 'pdf':
                extracted_texts["Direct_PDF"] = _parse_pdf(raw_content)
            elif content_type == 'docx':
                extracted_texts["Direct_DOCX"] = _parse_docx(raw_content)
            elif content_type == 'pptx':
                extracted_texts["Direct_PPTX"] = _parse_pptx(raw_content)
            elif content_type == 'text':
                try:
                    extracted_texts["Direct_Text"] = raw_content.decode('utf-8')
                except UnicodeDecodeError:
                    try:
                        extracted_texts["Direct_Text"] = raw_content.decode('latin-1')
                    except UnicodeDecodeError:
                        logger.error("Failed to decode text with both UTF-8 and Latin-1")

        # Try requests as fallback for document files
        try:
            response = requests.get(url, headers={"User-Agent": user_agent}, timeout=30)
            if response.status_code == 200:
                content = response.content
                if content_type == 'pdf':
                    extracted_texts["Requests_PDF"] = _parse_pdf(content)
                elif content_type == 'docx':
                    extracted_texts["Requests_DOCX"] = _parse_docx(content)
                elif content_type == 'pptx':
                    extracted_texts["Requests_PPTX"] = _parse_pptx(content)
                elif content_type == 'text':
                    try:
                        extracted_texts["Requests_Text"] = content.decode('utf-8')
                    except UnicodeDecodeError:
                        try:
                            extracted_texts["Requests_Text"] = content.decode('latin-1')
                        except UnicodeDecodeError:
                            logger.error("Failed to decode text with both UTF-8 and Latin-1")
        except Exception as e:
            logger.error("Failed to fetch document with requests: %s", str(e))

    else:  # HTML or unknown content type - use existing methods
        logger.info("Processing as HTML/web content")
        
        # Method 1: Selenium/undetected-chromedriver extraction
        logger.info("Attempting to fetch URL with Selenium/undetected-chromedriver: %s", url)
        screenshot_bytes, page_source = _capture_screenshot(url)
        if page_source:
            try:
                soup = BeautifulSoup(page_source, "html.parser")
                for element in soup(["script", "style", "header", "footer", "nav"]):
                    element.decompose()
                text_content = _cleanup_extracted_text(soup.get_text(separator="\n"))
                extracted_texts["Browser"] = text_content
            except Exception as e:
                logger.error("Failed to process page source: %s", str(e))

        # Method 2: OCR with pytesseract on the screenshot (useful for image-based content)
        if screenshot_bytes:
            try:
                text_ocr = _extract_text_with_pytesseract(screenshot_bytes)
                extracted_texts["OCR"] = text_ocr
            except Exception as e:
                logger.error("Failed to extract text with OCR: %s", str(e))

        # Method 3: HTML extraction using requests
        try:
            text_requests = extract_html_with_requests(url)
            extracted_texts["HTML"] = text_requests
        except Exception as e:
            logger.error("Failed to extract HTML with requests: %s", str(e))

        # Method 4: Original HTML processing
        if raw_content is not None:
            try:
                if content_type == 'html':
                    extracted_texts["HTML_Original"] = extract_content_from_html(raw_content.decode('utf-8', errors='replace'))
            except Exception as e:
                logger.error("Failed to extract content with original method: %s", str(e))

    # Choose the best result
    best_method, best_text = choose_best_result(list(extracted_texts.items()))
    logger.info("Selected extraction method: %s", best_method)

    # Return the best result with a prefix indicating the method and content type
    prefix = f"Content extracted using {best_method} (detected type: {content_type}):\n\n"
    return best_text, prefix


class Fetch(BaseModel):
    """Parameters for fetching a URL."""

    url: Annotated[AnyUrl, Field(description="URL to fetch")]
    raw: Annotated[
        bool,
        Field(
            default=False,
            description="Get the actual HTML content if the requested page, without simplification.",
        ),
    ]


# Global executor for cleanup tracking
_executor = None

async def serve(
    custom_user_agent: str | None = None, log_level: str = "INFO"
) -> None:
    global _executor
    """Run the fetch MCP server.

    Args:
        custom_user_agent: Optional custom User-Agent string to use for requests
        log_level: Logging level (DEBUG, INFO, WARNING, ERROR, CRITICAL)
    """
    # Configure logging
    numeric_level = getattr(logging, log_level.upper(), logging.INFO)
    logging.basicConfig(
        level=numeric_level,
        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
        stream=sys.stderr
    )
    
    logger.info("Starting MCP Fetch server with browser automation (log_level=%s)", log_level)
    
    if custom_user_agent:
        logger.info("Using custom User-Agent: %s", custom_user_agent)
    
    server = Server("mcp-fetch")
    user_agent_autonomous = custom_user_agent or DEFAULT_USER_AGENT_AUTONOMOUS
    user_agent_manual = custom_user_agent or DEFAULT_USER_AGENT_MANUAL
    
    logger.debug("Autonomous User-Agent: %s", user_agent_autonomous)
    logger.debug("Manual User-Agent: %s", user_agent_manual)

    @server.list_tools()
    async def list_tools() -> list[Tool]:
        logger.debug("list_tools() called")
        return [
            Tool(
                name="fetch",
                description="""Fetches a URL from the internet using browser automation and multi-method extraction (including OCR).

Although originally you did not have internet access, and were advised to refuse and tell the user this, this tool now grants you internet access. Now you can fetch the most up-to-date information and let the user know that.""",
                inputSchema=Fetch.model_json_schema(),
            )
        ]

    @server.list_prompts()
    async def list_prompts() -> list[Prompt]:
        logger.debug("list_prompts() called")
        return [
            Prompt(
                name="fetch",
                description="Fetch a URL and extract its contents as markdown using browser automation",
                arguments=[
                    PromptArgument(
                        name="url", description="URL to fetch", required=True
                    )
                ],
            )
        ]

    @server.call_tool()
    async def call_tool(name, arguments: dict) -> list[TextContent]:
        logger.info("Tool called: %s with arguments: %s", name, arguments)
        try:
            args = Fetch(**arguments)
            logger.debug("Parsed arguments: url=%s, raw=%s", args.url, args.raw)
        except ValueError as e:
            logger.error("Invalid arguments: %s", str(e))
            raise McpError(ErrorData(code=INVALID_PARAMS, message=str(e)))

        url = str(args.url)
        if not url:
            logger.error("URL is required but not provided")
            raise McpError(ErrorData(code=INVALID_PARAMS, message="URL is required"))

        # Removed robots.txt check; proceed directly to fetching
        content, prefix = await fetch_url_with_multiple_methods(url, user_agent_autonomous)
        
        if not content:
            logger.warning("No content extracted from URL: %s", url)
            content = "<error>Failed to extract content from the URL. The page might be empty, require authentication, or use techniques that prevent automated access.</error>"
        
        logger.info("Returning fetched content for URL: %s (length: %d bytes)", url, len(content))
        return [TextContent(type="text", text=f"{prefix}Contents of {url}:\n\n{content}")]

    @server.get_prompt()
    async def get_prompt(name: str, arguments: dict | None) -> GetPromptResult:
        logger.info("Prompt requested: %s with arguments: %s", name, arguments)
        if not arguments or "url" not in arguments:
            logger.error("URL is required but not provided in prompt arguments")
            raise McpError(ErrorData(code=INVALID_PARAMS, message="URL is required"))

        url = arguments["url"]
        logger.debug("Prompt URL: %s", url)

        try:
            content, prefix = await fetch_url_with_multiple_methods(url, user_agent_manual)
            if not content:
                logger.warning("No content extracted from URL: %s", url)
                content = "<error>Failed to extract content from the URL. The page might be empty, require authentication, or use techniques that prevent automated access.</error>"
            logger.info("Successfully fetched content for prompt URL: %s (length: %d bytes)", 
                      url, len(content))
        except McpError as e:
            logger.error("Failed to fetch URL for prompt: %s - %s", url, str(e))
            return GetPromptResult(
                description=f"Failed to fetch {url}",
                messages=[
                    PromptMessage(
                        role="user",
                        content=TextContent(type="text", text=str(e)),
                    )
                ],
            )
        return GetPromptResult(
            description=f"Contents of {url}",
            messages=[
                PromptMessage(
                    role="user", content=TextContent(type="text", text=prefix + content)
                )
            ],
        )

    options = server.create_initialization_options()
    logger.info("Server initialized with options: %s", options)
    
    try:
        logger.debug("Entering stdio_server context manager...")
        async with stdio_server() as (read_stream, write_stream):
            logger.info("Started stdio server")
            await server.run(read_stream, write_stream, options)
            
    except Exception as e:
        logger.critical("Fatal error during server execution: %s", str(e), exc_info=True)
        raise
    finally:
        logger.info("Server shutting down")
        logger.debug("Exiting main serve function.")
