from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os


def b_imgs2ppt(
        image_folder,
        output_ppt="output.pptx",
        sorted_key=lambda x: int(x.split(".")[0]),
        base_width_in=10  # 固定基准宽度（英寸）
):
    images = sorted([
        file for file in os.listdir(image_folder)
        if file.lower().endswith((".png", ".jpg", ".jpeg", ".bmp", ".gif"))
    ], key=sorted_key)

    prs = Presentation()

    for img_name in images:
        img_path = os.path.join(image_folder, img_name)

        with Image.open(img_path) as img:
            width_px, height_px = img.size
            aspect_ratio = height_px / width_px

        # 基于固定宽度，计算高度
        width_in = base_width_in
        height_in = base_width_in * aspect_ratio

        # 设置PPT页面大小
        prs.slide_width = Inches(width_in)
        prs.slide_height = Inches(height_in)

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(img_path, 0, 0,
                                 width=prs.slide_width, height=prs.slide_height)

    prs.save(output_ppt)
    print(f"PPT已保存到: {output_ppt}")
