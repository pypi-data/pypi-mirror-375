import os
from maximum_agents.base import Tool
from docx import Document
from pptx import Presentation

class GetPresentationTool(Tool):
    """Class-based tool for creating PowerPoint presentations"""
    
    name = "get_presentation"
    description = """
    Creates a new PowerPoint presentation with the given title and subtitle.
    
    Args:
        title (str): The title text to display on the first slide
        subtitle (str): The subtitle text to display on the first slide
        
    Returns:
        Presentation: A PowerPoint presentation object with initial title slide
    Creates and returns a PowerPoint presentation with various formatting examples.
    
    Example usage:
        # Hello World slide
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Hello, World!"
        subtitle.text = "python-pptx was here!"
        
        prs.save('test.pptx')

        # Bullet slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = 'Adding a Bullet Slide'
        
        tf = body_shape.text_frame
        tf.text = 'Find the bullet slide layout'
        
        p = tf.add_paragraph()
        p.text = 'Use _TextFrame.text for first bullet'
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
        p.level = 2

        # Add textbox
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        left = top = width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        tf.text = "This is text inside a textbox"
        
        p = tf.add_paragraph()
        p.text = "This is a second paragraph that's bold"
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "This is a third paragraph that's big"
        p.font.size = Pt(40)

        # Add picture
        img_path = 'monty-truth.png'
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        left = top = Inches(1)
        pic = slide.shapes.add_picture(img_path, left, top)
        
        left = Inches(5)
        height = Inches(5.5)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)

        # Add shapes
        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes
        
        shapes.title.text = 'Adding an AutoShape'
        
        left = Inches(0.93)
        top = Inches(3.0)
        width = Inches(1.75)
        height = Inches(1.0)
        
        shape = shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
        shape.text = 'Step 1'
        
        left = left + width - Inches(0.4)
        width = Inches(2.0)
        
        for n in range(2, 6):
            shape = shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
            shape.text = 'Step %d' % n
            left = left + width - Inches(0.4)

        # Add table
        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes
        
        shapes.title.text = 'Adding a Table'
        
        rows = cols = 2
        left = top = Inches(2.0)
        width = Inches(6.0)
        height = Inches(0.8)
        
        table = shapes.add_table(rows, cols, left, top, width, height).table
        
        table.columns[0].width = Inches(2.0)
        table.columns[1].width = Inches(4.0)
        
        table.cell(0, 0).text = 'Foo'
        table.cell(0, 1).text = 'Bar'
        table.cell(1, 0).text = 'Baz'
        table.cell(1, 1).text = 'Qux'

        # Extract text example
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
    """
    inputs = {}
    output_type = "object"
    
    def forward(self) -> Presentation:
        """
        Creates a new PowerPoint presentation.
        
        Returns:
            Presentation: A PowerPoint presentation object
        """
        prs = Presentation()
        return prs


class GetDocumentTool(Tool):
    """Class-based tool for creating Word documents"""
    
    name = "get_document"
    description = """
    Creates and returns a Word document with various formatting examples.
    
    Example usage:
        document = get_document()

        document.add_heading('Document Title', 0)

        p = document.add_paragraph('A plain paragraph having some ')
        p.add_run('bold').bold = True
        p.add_run(' and some ')
        p.add_run('italic.').italic = True

        document.add_heading('Heading, level 1', level=1)
        document.add_paragraph('Intense quote', style='Intense Quote')

        document.add_paragraph(
            'first item in unordered list', style='List Bullet'
        )
        document.add_paragraph(
            'first item in ordered list', style='List Number'
        )

        document.add_picture('monty-truth.png', width=Inches(1.25))

        records = (
            (3, '101', 'Spam'),
            (7, '422', 'Eggs'),
            (4, '631', 'Spam, spam, eggs, and spam')
        )

        table = document.add_table(rows=1, cols=3)
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = 'Qty'
        hdr_cells[1].text = 'Id'
        hdr_cells[2].text = 'Desc'
        for qty, id, desc in records:
            row_cells = table.add_row().cells
            row_cells[0].text = str(qty)
            row_cells[1].text = id
            row_cells[2].text = desc

        document.add_page_break()

        document.save('demo.docx')

        ## to insert images
        Candidate protocol
            >>> run = paragraph.add_run()
            >>> inline_shape = run.add_picture(file_like_image, MIME_type=None)
            >>> inline_shape.width = width
            >>> inline_shape.height = height
            Minimal XML
            This XML represents the working hypothesis of the minimum XML that must be inserted to add a working picture to a document:

            <pic:pic xmlns:pic="http://schemas.openxmlformats.org/drawingml/2006/picture">
            <pic:nvPicPr>
                <pic:cNvPr id="1" name="python-powered.png"/>
                <pic:cNvPicPr/>
            </pic:nvPicPr>
            <pic:blipFill>
                <a:blip r:embed="rId7"/>
                <a:stretch>
                <a:fillRect/>
                </a:stretch>
            </pic:blipFill>
            <pic:spPr>
                <a:xfrm>
                <a:off x="0" y="0"/>
                <a:ext cx="859536" cy="343814"/>
                </a:xfrm>
                <a:prstGeom prst="rect"/>
            </pic:spPr>
            </pic:pic>

    Returns:
        Document: A python-docx Document object
    """
    inputs = {}
    output_type = "object"
    
    def forward(self) -> Document:
        """
        Creates a new Word document.
        
        Returns:
            Document: A python-docx Document object
        """
        document = Document()
        return document


class GetClientTool(Tool):
    """Class-based tool for getting Salesforce client"""
    
    name = "get_client"
    description = """
    Query Salesforce using SOQL.
     Returns:
        Salesforce client instance ready for API calls
        
    Usage Examples:
        # Get the client
        sf = get_client()
        
        # Query records
        accounts = sf.query("SELECT Id, Name FROM Account LIMIT 10")
        contacts = sf.query("SELECT Id, Email FROM Contact WHERE LastName = 'Smith'")
        Example:

        sf.query("SELECT Id, FirstName, LastName FROM Contact WHERE FirstName='Felix'")
        Return
        [{
            'attributes': {
                'type': 'Contact',
                'url': '/services/data/v44.0/sobjects/Contact/0031l000007Jia4AAC'
            },
            'Id': '0031l000007Jia4AAC',
            'FirstName': 'Felix',
            'LastName': 'Lindstrom'
        }, ...]

        The query function returns a list of dictionaries.
        
        # Create records
        new_account = sf.Account.insert({"Name": "Test Account", "Type": "Customer"})
        
        # Update records
        sf.Account.update("003XX000004TMM2", {"Name": "Updated Name"})
        
        # Delete records
        sf.Account.delete("003XX000004TMM2")

        sf.Contact.insert({"LastName": "Smith", "Email": "smith@example.com"})
        
    Common Salesforce Objects:
        - Account: Customer/company records
        - Contact: Individual person records  
        - Opportunity: Sales opportunities
        - Lead: Potential customers
        - Case: Customer service cases
        - User: Salesforce users
        - Custom Objects: Your org's custom objects (usually end with __c)
        
    SOQL Query Examples:
        - "SELECT Id, Name, Type FROM Account WHERE Type = 'Customer'"
        - "SELECT Id, FirstName, LastName, Email FROM Contact WHERE Account.Name = 'Acme Corp'"
        - "SELECT Id, Name, Amount, StageName FROM Opportunity WHERE CloseDate = THIS_MONTH"
    """
    inputs = {}
    output_type = "object"
    
    def forward(self):
        """
        Get Salesforce client using environment variables.
        
        Returns:
            Salesforce client instance ready for API calls
        """
        from salesforce_api import Salesforce
        
        username = os.getenv('SALESFORCE_USERNAME')
        password = os.getenv('SALESFORCE_PASSWORD')
        security_token = os.getenv('SALESFORCE_TOKEN')
        
        if not all([username, password, security_token]):
            raise ValueError("Missing required environment variables: SALESFORCE_USERNAME, SALESFORCE_PASSWORD, SALESFORCE_TOKEN")
        
        # Type assertions since we've checked they're not None
        assert username is not None
        assert password is not None
        assert security_token is not None
        
        client = Salesforce(
            username=username,
            password=password,
            security_token=security_token,
        )
        
        return client.sobjects

